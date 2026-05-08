#!/usr/bin/env python3
"""
slides-en.py — parchment design system, English slide deck generator.

Usage:
  pip install python-pptx --break-system-packages
  python3 slides-en.py

Output:
  output.pptx (16:9, parchment aesthetic, Inter + Newsreader fonts)

This is a template. Fill in your content and run it directly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design system constants
# ═══════════════════════════════════════════════════════════

PARCHMENT  = RGBColor(0xf5, 0xf4, 0xed)
IVORY      = RGBColor(0xfa, 0xf9, 0xf5)
BRAND      = RGBColor(0x1B, 0x36, 0x5D)
BRAND_DEEP = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK = RGBColor(0x14, 0x14, 0x13)
DARK_WARM  = RGBColor(0x3d, 0x3d, 0x3a)
CHARCOAL   = RGBColor(0x4d, 0x4c, 0x48)
OLIVE      = RGBColor(0x5e, 0x5d, 0x59)
STONE      = RGBColor(0x87, 0x86, 0x7f)
BORDER     = RGBColor(0xe8, 0xe6, 0xdc)
WHITE      = RGBColor(0xff, 0xff, 0xff)

# English Silicon Valley stack. PowerPoint falls back silently if the
# primary face is not installed on the viewing machine.
SERIF = "Newsreader"    # fallback: Charter -> Georgia
SANS  = "Inter"         # fallback: Helvetica Neue -> Arial

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_weight=0.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


# ═══════════════════════════════════════════════════════════
# Slide templates
# ═══════════════════════════════════════════════════════════

def cover_slide(prs, title, subtitle, author, date):
    s = blank_slide(prs)
    add_text(s, title,
             Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             font=SERIF, size=48, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.3), Inches(1), weight_pt=1.5)
    add_text(s, subtitle,
             Inches(1), Inches(4.6), Inches(11.33), Inches(0.8),
             font=SANS, size=18, italic=True, color=OLIVE,
             align=PP_ALIGN.CENTER)
    add_text(s, f"{author} · {date}",
             Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE,
             align=PP_ALIGN.CENTER)
    return s


def toc_slide(prs, items):
    s = blank_slide(prs)
    add_text(s, "Contents",
             Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             font=SERIF, size=34, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)

    for i, item in enumerate(items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}",
                 Inches(1.2), y, Inches(1), Inches(0.6),
                 font=SERIF, size=28, color=BRAND)
        add_text(s, item,
                 Inches(2.4), y, Inches(9), Inches(0.6),
                 font=SERIF, size=22, color=NEAR_BLACK,
                 vanchor=MSO_ANCHOR.MIDDLE)
    return s


def chapter_slide(prs, number, title):
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, f"0{number}",
             Inches(0.8), Inches(0.5), Inches(2), Inches(0.8),
             font=SERIF, size=28, color=WHITE)
    add_text(s, title,
             Inches(1), Inches(3), Inches(11.33), Inches(1.5),
             font=SERIF, size=60, color=WHITE,
             align=PP_ALIGN.CENTER)
    return s


def content_slide(prs, eyebrow, title, body, page_num=None):
    s = blank_slide(prs)
    add_text(s, eyebrow.upper(),
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=11, color=STONE)
    add_text(s, title,
             Inches(1.2), Inches(1.2), Inches(11.33), Inches(1.2),
             font=SERIF, size=34, color=NEAR_BLACK)
    add_text(s, body,
             Inches(1.2), Inches(3), Inches(11), Inches(3.5),
             font=SANS, size=18, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f"— {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def metrics_slide(prs, title, metrics):
    """metrics: [(value, label), ...]"""
    s = blank_slide(prs)
    add_text(s, title,
             Inches(1.2), Inches(0.8), Inches(11), Inches(1),
             font=SERIF, size=30, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(2), Inches(1))

    n = len(metrics)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    start = (SLIDE_W - total_w) / 2

    for i, (value, label) in enumerate(metrics):
        x = start + (card_w + gap) * i
        add_text(s, value,
                 x, Inches(3), card_w, Inches(1.5),
                 font=SERIF, size=56, color=BRAND,
                 align=PP_ALIGN.CENTER)
        add_text(s, label,
                 x, Inches(4.8), card_w, Inches(0.6),
                 font=SANS, size=14, color=OLIVE,
                 align=PP_ALIGN.CENTER)
    return s


def quote_slide(prs, quote, source):
    s = blank_slide(prs)
    add_text(s, f"\u201c{quote}\u201d",
             Inches(1.5), Inches(2.8), Inches(10.33), Inches(2.5),
             font=SERIF, size=30, italic=True, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER,
             vanchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f"— {source}",
             Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.4),
             font=SANS, size=14, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


def ending_slide(prs, message, contact):
    s = blank_slide(prs)
    add_text(s, message,
             Inches(1), Inches(3), Inches(11.33), Inches(1.2),
             font=SERIF, size=44, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, contact,
             Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
             font=SANS, size=16, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


# ═══════════════════════════════════════════════════════════
# Main — example deck, replace with your content
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs,
        title="{{DOCUMENT_TITLE}}",
        subtitle="{{One-line description.}}",
        author="{{AUTHOR}}",
        date="2026.04")

    toc_slide(prs, items=[
        "{{Chapter 1}}",
        "{{Chapter 2}}",
        "{{Chapter 3}}",
        "Q&A",
    ])

    chapter_slide(prs, 1, "{{Chapter Title}}")

    content_slide(prs,
        eyebrow="{{Chapter · This page}}",
        title="{{Core claim as a sentence}}",
        body=("{{A short body paragraph, 18pt sans. Keep it under three lines. "
              "One slide, one core idea. The reader's attention is the scarce resource.}}"),
        page_num=4)

    metrics_slide(prs,
        title="Key Results",
        metrics=[
            ("+42%",   "Conversion lift"),
            ("3.8M",   "Monthly actives"),
            ("99.9%",  "Availability SLA"),
            ("5,000+", "QPS peak"),
        ])

    quote_slide(prs,
        quote="Good design is as little design as possible.",
        source="Dieter Rams")

    ending_slide(prs,
        message="Thank you",
        contact="{{EMAIL}} · {{WEBSITE}}")

    prs.save('output.pptx')
    print("✓ Saved output.pptx")


if __name__ == '__main__':
    main()
