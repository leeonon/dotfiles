#!/usr/bin/env python3
"""
gen_slides.py - parchment design system slide deck generator

用法：
  pip install python-pptx --break-system-packages
  python3 gen_slides.py

输出：
  output.pptx (16:9 宽屏, parchment 风格)

这是一个模板脚本 - 填充自己的内容后直接运行。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design System Constants
# ═══════════════════════════════════════════════════════════

# 色板
PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
BRAND_DEEP  = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
CHARCOAL    = RGBColor(0x4d, 0x4c, 0x48)
OLIVE       = RGBColor(0x5e, 0x5d, 0x59)
STONE       = RGBColor(0x87, 0x86, 0x7f)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

# 字体（PPT 会按序 fallback）
SERIF = "Source Han Serif SC"   # 中文 serif
SANS  = "Source Han Sans SC"    # 中文 sans

# 16:9 宽屏
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg_color=PARCHMENT):
    """创建空白幻灯片，指定背景色"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # 6 = Blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    """加一段文字"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    """加水平线"""
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_weight=0.5):
    """加卡片背景"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


# ═══════════════════════════════════════════════════════════
# Slide Templates
# ═══════════════════════════════════════════════════════════

def cover_slide(prs, title, subtitle, author, date):
    """封面：大标题 + 副标题 + 作者/日期"""
    s = blank_slide(prs)
    # 大标题（serif 44pt 居中）
    add_text(s, title,
             Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             font=SERIF, size=44, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    # 品牌色短线
    add_line(s, Inches(6.17), Inches(4.2), Inches(1), weight_pt=1.5)
    # 副标题
    add_text(s, subtitle,
             Inches(1), Inches(4.5), Inches(11.33), Inches(0.8),
             font=SANS, size=18, color=OLIVE,
             align=PP_ALIGN.CENTER)
    # 作者 + 日期
    add_text(s, f"{author}　·　{date}",
             Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE,
             align=PP_ALIGN.CENTER)
    return s


def toc_slide(prs, items):
    """目录页：01 章节名 列表"""
    s = blank_slide(prs)
    add_text(s, "目录",
             Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)

    for i, item in enumerate(items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}",
                 Inches(1.2), y, Inches(1), Inches(0.6),
                 font=SERIF, size=28, color=BRAND)
        add_text(s, item,
                 Inches(2.4), y, Inches(9), Inches(0.6),
                 font=SERIF, size=22, color=NEAR_BLACK,
                 vanchor=MSO_ANCHOR.MIDDLE)
    return s


def chapter_slide(prs, number, title):
    """章节首页：油墨蓝色背景 + 居中大标题"""
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, f"0{number}",
             Inches(0.8), Inches(0.5), Inches(2), Inches(0.8),
             font=SERIF, size=26, color=WHITE)
    add_text(s, title,
             Inches(1), Inches(3), Inches(11.33), Inches(1.5),
             font=SERIF, size=56, color=WHITE,
             align=PP_ALIGN.CENTER)
    return s


def content_slide(prs, eyebrow, title, body, page_num=None):
    """内容页：小标题 + 大标题 + 正文"""
    s = blank_slide(prs)
    # eyebrow
    add_text(s, eyebrow,
             Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    # title
    add_text(s, title,
             Inches(1.2), Inches(1.2), Inches(11.33), Inches(1.2),
             font=SERIF, size=32, color=NEAR_BLACK)
    # body
    add_text(s, body,
             Inches(1.2), Inches(3), Inches(11), Inches(3.5),
             font=SANS, size=18, color=DARK_WARM)
    # page number
    if page_num is not None:
        add_text(s, f"— {page_num:02d}",
                 Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE,
                 align=PP_ALIGN.RIGHT)
    return s


def metrics_slide(prs, title, metrics):
    """数据页：标题 + N 张数据卡并排
    metrics: [(value, label), ...]
    """
    s = blank_slide(prs)
    # 标题
    add_text(s, title,
             Inches(1.2), Inches(0.8), Inches(11), Inches(1),
             font=SERIF, size=28, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(2), Inches(1))

    # 数据卡
    n = len(metrics)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    start = (SLIDE_W - total_w) / 2

    for i, (value, label) in enumerate(metrics):
        x = start + (card_w + gap) * i
        # 大数字
        add_text(s, value,
                 x, Inches(3), card_w, Inches(1.5),
                 font=SERIF, size=52, color=BRAND,
                 align=PP_ALIGN.CENTER)
        # 标签
        add_text(s, label,
                 x, Inches(4.8), card_w, Inches(0.6),
                 font=SANS, size=14, color=OLIVE,
                 align=PP_ALIGN.CENTER)
    return s


def quote_slide(prs, quote, source):
    """引用页：极简，居中斜体引文"""
    s = blank_slide(prs)
    add_text(s, f"\u201c{quote}\u201d",
             Inches(1.5), Inches(2.8), Inches(10.33), Inches(2.5),
             font=SERIF, size=28, italic=True, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER,
             vanchor=MSO_ANCHOR.MIDDLE)
    add_text(s, f"— {source}",
             Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.4),
             font=SANS, size=14, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


def ending_slide(prs, message, contact):
    """结束页"""
    s = blank_slide(prs)
    add_text(s, message,
             Inches(1), Inches(3), Inches(11.33), Inches(1.2),
             font=SERIF, size=40, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, contact,
             Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
             font=SANS, size=16, color=OLIVE,
             align=PP_ALIGN.CENTER)
    return s


# ═══════════════════════════════════════════════════════════
# Main: 示例 deck，按实际需求改
# ═══════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. 封面
    cover_slide(prs,
        title="{{文档标题}}",
        subtitle="{{一句话描述}}",
        author="{{作者}}",
        date="2026.04")

    # 2. 目录
    toc_slide(prs, items=[
        "{{章节 1}}",
        "{{章节 2}}",
        "{{章节 3}}",
        "{{Q&A}}",
    ])

    # 3. 章节首页
    chapter_slide(prs, 1, "{{章节标题}}")

    # 4. 内容页
    content_slide(prs,
        eyebrow="{{章节 · 本页}}",
        title="{{核心论点标题}}",
        body="{{一段正文，18pt sans 字体。控制在 3 行内，一屏一个核心信息。}}",
        page_num=4)

    # 5. 数据页
    metrics_slide(prs,
        title="关键结果",
        metrics=[
            ("+42%",   "转化率提升"),
            ("3.8M",   "月活用户"),
            ("99.9%",  "可用性 SLA"),
            ("5,000+", "QPS 峰值"),
        ])

    # 6. 引用
    quote_slide(prs,
        quote="好的设计是尽可能少的设计。",
        source="Dieter Rams")

    # 7. 结束
    ending_slide(prs,
        message="Thank you",
        contact="{{邮箱}}　·　{{网站}}")

    prs.save('output.pptx')
    print("✓ Saved output.pptx")


if __name__ == '__main__':
    main()
